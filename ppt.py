from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Blockchain-Based Voting System"
subtitle.text = "A Secure, Transparent, and Tamper-Proof Electronic Voting Platform\n\nCapstone Project"

# Slide 2: Agenda
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Agenda"
content = slide.placeholders[1]
content.text = (
    "• Problem Statement\n"
    "• Objectives\n"
    "• What is Blockchain?\n"
    "• System Architecture\n"
    "• Core Features\n"
    "• Security Enhancements\n"
    "• Workflow\n"
    "• Implementation Details\n"
    "• Demonstration\n"
    "• Results & Testing\n"
    "• Conclusion & Future Work\n"
    "• Q&A"
)

# Slide 3: Problem Statement
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Problem Statement"
content = slide.placeholders[1]
content.text = (
    "Traditional Voting Systems Face Issues:\n\n"
    "• Voter fraud – impersonation, double voting\n"
    "• Lack of transparency – results cannot be independently verified\n"
    "• Tampering – centralized databases can be altered\n"
    "• No voter verifiability – voters cannot confirm their vote was counted\n"
    "• High cost – printing ballots, polling stations, manual counting\n\n"
    "Need for a secure, decentralized, and transparent voting solution."
)

# Slide 4: Objectives
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Project Objectives"
content = slide.placeholders[1]
content.text = (
    "• Build a blockchain-based voting system where every vote is a transaction\n"
    "• Ensure immutability – votes cannot be changed once recorded\n"
    "• Prevent double voting using cryptographic voter tracking\n"
    "• Provide end-to-end verifiability – anyone can audit the blockchain\n"
    "• Implement digital signatures to authenticate voters\n"
    "• Allow admin controls (mining, reset) and CSV export of results"
)

# Slide 5: What is Blockchain?
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "What is Blockchain?"
content = slide.placeholders[1]
content.text = (
    "Definition: A distributed, decentralized ledger that records transactions in blocks.\n\n"
    "Key Properties:\n"
    "• Immutability – data cannot be altered without changing all subsequent blocks\n"
    "• Transparency – anyone can view the entire chain\n"
    "• Decentralization – no single point of failure\n"
    "• Consensus – proof-of-work ensures agreement\n\n"
    "Each block contains: index, timestamp, transactions, previous hash, nonce, its own hash"
)

# Slide 6: System Architecture
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "System Architecture"
content = slide.placeholders[1]
content.text = (
    "Components:\n"
    "• Frontend: HTML/CSS/JS, Web Crypto API for key generation & signing\n"
    "• Backend: Flask (Python) – handles HTTP requests, blockchain logic\n"
    "• Blockchain module: Block class, Blockchain class, proof-of-work, persistence\n"
    "• Storage: JSON files – blockchain.json and eligible_voters.json\n\n"
    "Data Flow: Browser ↔ Flask Server ↔ Blockchain Core ↔ JSON files"
)

# Slide 7: Core Features (1)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Core Features (1)"
content = slide.placeholders[1]
content.text = (
    "• Vote casting – Voter enters ID, selects candidate, vote is signed and sent\n"
    "• Pending transactions – Votes wait in a pool until mined\n"
    "• Mining (proof-of-work) – Admin collects pending votes into a new block\n"
    "• Blockchain explorer – Anyone can view full chain and verify integrity\n"
    "• Result tally – Counts votes from all blocks in the chain\n"
    "• Persistent storage – Blockchain saved to blockchain.json – survives restarts"
)

# Slide 8: Core Features (2) – Security Enhancements
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Security Enhancements"
content = slide.placeholders[1]
content.text = (
    "• Digital signatures – RSA-2048, SHA-256. Voter signs 'voterID:candidate'\n"
    "• Voter eligibility – Admin uploads CSV of allowed voter IDs\n"
    "• Double voting prevention – Set of voted IDs tracked across chain and pending pool\n"
    "• Admin-only mining – Password-protected mining endpoint\n"
    "• Admin reset – Delete all votes, pending transactions, eligible list\n"
    "• CSV export – Download results as CSV for external auditing"
)

# Slide 9: Workflow
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Workflow – Step by Step"
content = slide.placeholders[1]
content.text = (
    "1. Admin uploads eligible voters (CSV)\n"
    "2. Voter enters ID and candidate\n"
    "3. Browser generates RSA key pair (first time) or retrieves from localStorage\n"
    "4. Browser signs 'voterID:candidate' using private key\n"
    "5. Browser sends vote (ID, candidate, signature, public key) to /vote\n"
    "6. Server verifies eligibility, no double vote, and signature\n"
    "7. If valid, vote added to pending_transactions\n"
    "8. Admin clicks 'Mine Pending Votes'\n"
    "9. Server creates new block, solves proof-of-work, appends to chain\n"
    "10. Results updated – anyone can refresh or export CSV"
)

# Slide 10: Digital Signatures
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Digital Signatures"
content = slide.placeholders[1]
content.text = (
    "Key Generation (in browser):\n"
    "• Web Crypto API generates RSA-2048 key pair\n"
    "• Private key stored in browser's localStorage (never sent to server)\n"
    "• Public key sent to server with each vote\n\n"
    "Signing:\n"
    "message = `${voterId}:${candidate}`;\n"
    "signature = await crypto.subtle.sign(...);\n\n"
    "Verification (on server):\n"
    "message = f'{voter_id}:{candidate}'.encode()\n"
    "pub_key.verify(signature, message)\n\n"
    "Only the owner of the private key can vote for that ID."
)

# Slide 11: Proof-of-Work (Mining)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Proof-of-Work (Mining)"
content = slide.placeholders[1]
content.text = (
    "Purpose: Prevent tampering – changing a past vote would require re-mining all later blocks.\n\n"
    "Algorithm:\n"
    "• Find a nonce such that hash(block + nonce) starts with difficulty number of zeros\n"
    "• Difficulty = 4 → target '0000'\n\n"
    "Code snippet:\n"
    "def mine_block(self, difficulty):\n"
    "    target = '0' * difficulty\n"
    "    while self.hash[:difficulty] != target:\n"
    "        self.nonce += 1\n"
    "        self.hash = self.compute_hash()\n\n"
    "Performance: ~0.5‑2 seconds per block on a modern CPU."
)

# Slide 12: Data Persistence & APIs
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Persistence & APIs"
content = slide.placeholders[1]
content.text = (
    "Files created automatically:\n"
    "• blockchain.json – full chain, pending transactions, voted set\n"
    "• eligible_voters.json – list of allowed voter IDs\n\n"
    "Key REST API endpoints:\n"
    "• POST /vote – cast a signed vote\n"
    "• POST /mine – mine pending votes (admin password)\n"
    "• POST /admin/upload_eligible – upload CSV of eligible voters\n"
    "• POST /admin/reset – delete all data (admin password)\n"
    "• GET /results – vote counts as JSON\n"
    "• GET /export_results – download CSV of results\n"
    "• GET /chain – view full blockchain"
)

# Slide 13: Screenshots (placeholder)
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Screenshots / Demo"
content = slide.placeholders[1]
content.text = (
    "Insert actual screenshots here:\n\n"
    "1. Home page – voting form, admin sections\n"
    "2. Upload CSV – success message\n"
    "3. Casting a vote – 'Vote recorded (pending mining)'\n"
    "4. Mining – 'Block 1 mined'\n"
    "5. Results table – vote counts\n"
    "6. Export CSV – file downloaded\n"
    "7. Blockchain explorer – JSON view of blocks\n"
    "8. Reset confirmation – all data cleared"
)

# Slide 14: Testing & Validation
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Testing & Validation"
content = slide.placeholders[1]
content.text = (
    "• Vote with eligible ID → accepted (pending)\n"
    "• Vote with ineligible ID → 'Not eligible' error\n"
    "• Double vote same ID → 'Already voted' error\n"
    "• Tamper with a block (manually edit JSON) → is_chain_valid() returns false\n"
    "• Mine with wrong password → 'Unauthorized'\n"
    "• Reset system → chain resets to genesis block, eligible list empty\n"
    "• Signature mismatch → verification fails"
)

# Slide 15: Advantages
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Advantages of Our System"
content = slide.placeholders[1]
content.text = (
    "✅ Transparent – anyone can view the full blockchain\n"
    "✅ Tamper-proof – proof-of-work makes altering past votes expensive\n"
    "✅ No double voting – voter IDs tracked immutably\n"
    "✅ Voter authentication – digital signatures prevent impersonation\n"
    "✅ Admin control – mining and reset protected by password\n"
    "✅ Portable – runs on any machine with Python\n"
    "✅ Easy to audit – CSV export and JSON chain explorer"
)

# Slide 16: Limitations & Future Work
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Limitations & Future Work"
content = slide.placeholders[1]
content.text = (
    "Current limitations:\n"
    "• Centralized server (not fully decentralized)\n"
    "• No voter anonymity (voter ID is visible)\n"
    "• Proof-of-work consumes CPU (fine for small scale)\n"
    "• Admin password is hardcoded\n\n"
    "Future enhancements:\n"
    "• Peer-to-peer network with consensus (multiple nodes)\n"
    "• Zero-knowledge proofs or ring signatures for anonymity\n"
    "• Smart contract version on Ethereum (Solidity)\n"
    "• Voter receipt generation (end-to-end verifiability)\n"
    "• Docker containerisation for easy deployment"
)

# Slide 17: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
content.text = (
    "Successfully built a complete blockchain voting system with:\n"
    "• Immutable vote ledger\n"
    "• Digital signatures for authentication\n"
    "• Voter eligibility list\n"
    "• Admin mining and reset\n"
    "• CSV export\n\n"
    "Demonstrates core blockchain concepts: hashing, proof-of-work, immutability.\n"
    "Provides a practical, working prototype suitable for small-scale elections.\n"
    "Can be extended to a fully decentralized system.\n\n"
    "Thank you!"
)

# Slide 18: Q&A
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Q&A"
content = slide.placeholders[1]
content.text = "Questions?"

# Save the presentation
prs.save("Blockchain_Voting_System.pptx")
print("Presentation created: Blockchain_Voting_System.pptx")